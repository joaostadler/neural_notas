"""
app/routes/main.py - Rotas principais da interface web.
"""

import os
from datetime import date, datetime, timedelta
from pathlib import Path
from uuid import uuid4

from flask import Blueprint, current_app, flash, jsonify, make_response, redirect, render_template, request, send_file, url_for
from flask_login import current_user, login_required
from sqlalchemy import desc
from werkzeug.utils import secure_filename

from models import (
    Apresentacao, Caderno, CartaoKanban, CelulaJupyter, ColunaKanban,
    Codigo, ComentarioSlide, Diagrama, HistoricoEtapas, IconeCustomizado,
    Imagem, Nota, PDF, Planilha, TarefaFacil, Topico, db,
)

bp_main = Blueprint('main', __name__)


@bp_main.route('/health')
def health():
    """Keep-alive endpoint — sem autenticação."""
    from flask import jsonify
    from datetime import datetime, timezone, timedelta
    brt = timezone(timedelta(hours=-3))
    return jsonify({
        'status': 'ok',
        'ts': datetime.now(brt).strftime('%Y-%m-%d %H:%M:%S BRT'),
    })


TIPOS_DOCUMENTO = {
    'pasta':    {'icone': '🗂',  'label': 'Pasta'},
    'nota':     {'icone': '📒', 'label': 'Nota'},
    'caderno':  {'icone': '📔', 'label': 'Caderno'},
    'planilha': {'icone': '📊', 'label': 'Planilha'},
    'diagrama': {'icone': '🔷', 'label': 'Diagrama'},
    'jupyter':  {'icone': '⚗️',  'label': 'Notebook'},
    'tarefa':   {'icone': '🎯', 'label': 'Tarefa'},
    'imagem':   {'icone': '🖼',  'label': 'Imagem'},
    'python':   {'icone': 'PY',  'label': 'Python'},
    'sql':      {'icone': 'SQL', 'label': 'SQL'},
    'biblioteca':    {'icone': '📚', 'label': 'Biblioteca'},
    'apresentacao':  {'icone': '📊', 'label': 'Apresentação'},
}


def _topico_do_usuario(topico_id):
    if not topico_id:
        return None
    return Topico.query.filter_by(id=topico_id, id_usuario=current_user.id).first()


def _eh_container(topico):
    return topico and topico.tipo in ('pasta', 'biblioteca')


def _proxima_ordem(id_pai):
    ultima_ordem = (
        db.session.query(db.func.max(Topico.ordem))
        .filter_by(id_usuario=current_user.id, id_pai=id_pai)
        .scalar()
    )
    return (ultima_ordem or 0) + 1


def _criar_conteudo_inicial(topico, tipo):
    if tipo in ('nota', 'tarefa'):
        db.session.add(Nota(id_topico=topico.id, titulo=topico.nome, conteudo=''))
    elif tipo == 'caderno':
        db.session.add(Caderno(id_topico=topico.id, titulo=topico.nome, conteudo=''))
    elif tipo == 'planilha':
        db.session.add(Planilha(id_topico=topico.id, titulo=topico.nome, dados_json='{"colunas":["A","B","C"],"linhas":[["","",""],["","",""],["","",""]]}'))
    elif tipo == 'jupyter':
        db.session.add(CelulaJupyter(id_topico=topico.id, tipo='markdown', conteudo=f'# {topico.nome}', ordem=1))
    elif tipo == 'diagrama':
        db.session.add(Diagrama(id_topico=topico.id, titulo=topico.nome, dados_json=''))
    elif tipo == 'python':
        db.session.add(Codigo(id_topico=topico.id, titulo=topico.nome, conteudo='# Python\n'))
    elif tipo == 'sql':
        db.session.add(Codigo(id_topico=topico.id, titulo=topico.nome, conteudo='-- SQL\n'))


def _caminho_static_seguro(caminho_relativo):
    """Converte um caminho salvo no banco em arquivo dentro de app/static."""
    if not caminho_relativo:
        return None

    static_root = Path(current_app.static_folder).resolve()
    arquivo = (static_root / caminho_relativo).resolve()
    try:
        arquivo.relative_to(static_root)
    except ValueError:
        current_app.logger.warning('Ignorando caminho fora da pasta static: %s', caminho_relativo)
        return None
    return arquivo


def _coletar_arquivos_do_topico(topico):
    arquivos = []

    if topico.pdfs and topico.pdfs.caminho:
        arquivos.append(topico.pdfs.caminho)
    if topico.imagens and topico.imagens.caminho:
        arquivos.append(topico.imagens.caminho)

    for subtopico in topico.subtopicos:
        arquivos.extend(_coletar_arquivos_do_topico(subtopico))

    return arquivos


def _remover_arquivos_static(caminhos_relativos):
    uploads_root = Path(current_app.config['UPLOAD_FOLDER']).resolve()

    for caminho_relativo in caminhos_relativos:
        arquivo = _caminho_static_seguro(caminho_relativo)
        if not arquivo or not arquivo.is_file():
            continue

        try:
            arquivo.unlink()
            pasta = arquivo.parent
            while pasta != uploads_root and uploads_root in pasta.parents:
                try:
                    pasta.rmdir()
                except OSError:
                    break
                pasta = pasta.parent
        except OSError as exc:
            current_app.logger.warning('Nao foi possivel remover arquivo %s: %s', arquivo, exc)


@bp_main.route('/')
def index():
    """Pagina inicial."""
    if current_user.is_authenticated:
        return redirect(url_for('main.dashboard'))
    return redirect(url_for('auth.login'))


@bp_main.route('/dashboard')
@login_required
def dashboard():
    """Dashboard principal do usuario."""
    topicos = (
        Topico.query.filter_by(id_usuario=current_user.id, id_pai=None)
        .order_by(Topico.ordem, Topico.nome)
        .all()
    )

    tarefas_do_dia = (
        TarefaFacil.query
        .filter(
            TarefaFacil.id_usuario == current_user.id,
            TarefaFacil.data_tarefa == date.today(),
        )
        .order_by(TarefaFacil.concluida.asc(), TarefaFacil.criado_em.asc())
        .limit(15)
        .all()
    )

    icones_customizados = {
        ic.tipo: f'/static/{ic.caminho}'
        for ic in IconeCustomizado.query.filter_by(id_usuario=current_user.id).all()
    }

    return render_template(
        'main/dashboard.html',
        topicos=topicos,
        tipos_documento=TIPOS_DOCUMENTO,
        tarefas_do_dia=tarefas_do_dia,
        icones_customizados=icones_customizados,
    )


@bp_main.route('/topicos/criar', methods=['POST'])
@login_required
def criar_topico():
    """Cria pastas e documentos dentro da arvore lateral."""
    nome = request.form.get('nome', '').strip()
    tipo = request.form.get('tipo', '').strip()
    id_pai_raw = request.form.get('id_pai', '').strip()
    id_pai = int(id_pai_raw) if id_pai_raw.isdigit() else None

    if not nome:
        flash('Informe um nome para criar o item.', 'erro')
        return redirect(url_for('main.dashboard'))

    if tipo not in TIPOS_DOCUMENTO:
        flash('Tipo de item invalido.', 'erro')
        return redirect(url_for('main.dashboard'))

    if id_pai and not _eh_container(_topico_do_usuario(id_pai)):
        flash('A pasta de destino nao foi encontrada.', 'erro')
        return redirect(url_for('main.dashboard'))

    topico = Topico(
        id_usuario=current_user.id,
        id_pai=id_pai,
        nome=nome,
        tipo=tipo,
        icone=TIPOS_DOCUMENTO[tipo]['icone'],
        ordem=_proxima_ordem(id_pai),
    )
    db.session.add(topico)
    db.session.flush()  # Garante que o topico tem um ID antes de criar o conteúdo
    _criar_conteudo_inicial(topico, tipo)
    db.session.commit()

    flash(f'{TIPOS_DOCUMENTO[tipo]["label"]} criado com sucesso.', 'sucesso')
    return redirect(url_for('main.dashboard'))


@bp_main.route('/pdf/importar', methods=['POST'])
@login_required
def importar_pdf():
    """Importa um PDF como item da arvore lateral."""
    arquivo = request.files.get('arquivo_pdf')
    id_pai_raw = request.form.get('id_pai_pdf', '').strip()
    id_pai = int(id_pai_raw) if id_pai_raw.isdigit() else None

    if id_pai and not _eh_container(_topico_do_usuario(id_pai)):
        flash('A pasta de destino nao foi encontrada.', 'erro')
        return redirect(url_for('main.dashboard'))

    if not arquivo or not arquivo.filename:
        flash('Escolha um arquivo PDF para importar.', 'erro')
        return redirect(url_for('main.dashboard'))

    nome_seguro = secure_filename(arquivo.filename)
    if not nome_seguro.lower().endswith('.pdf'):
        flash('Somente arquivos PDF podem ser importados aqui.', 'erro')
        return redirect(url_for('main.dashboard'))

    conteudo_bytes = arquivo.read()
    titulo = os.path.splitext(nome_seguro)[0] or 'PDF importado'
    topico = Topico(
        id_usuario=current_user.id,
        id_pai=id_pai,
        nome=titulo,
        tipo='pdf',
        icone='PDF',
        ordem=_proxima_ordem(id_pai),
    )
    db.session.add(topico)
    db.session.flush()
    db.session.add(PDF(id_topico=topico.id, titulo=titulo, conteudo_pdf=conteudo_bytes))
    db.session.commit()

    flash('PDF importado com sucesso.', 'sucesso')
    return redirect(url_for('main.dashboard'))


# ── Apresentações PPT/PPTX ────────────────────────────────────────────────────

def _extrair_slides_pptx(conteudo_bytes):
    """Extrai título, texto e imagens de cada slide via python-pptx."""
    import base64
    import json
    from io import BytesIO
    from pptx import Presentation

    prs = Presentation(BytesIO(conteudo_bytes))
    slides = []
    for slide in prs.slides:
        s = {'titulo': '', 'linhas': [], 'imagens': [], 'notas': ''}

        try:
            if slide.has_notes_slide:
                nf = slide.notes_slide.notes_text_frame
                if nf:
                    s['notas'] = nf.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            try:
                ph_idx = None
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    ph_idx = shape.placeholder_format.idx

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if not t:
                            continue
                        if ph_idx == 0:
                            if not s['titulo']:
                                s['titulo'] = t
                        else:
                            s['linhas'].append(t)

                # Detecta imagem pelo atributo image (mais compatível que MSO_SHAPE_TYPE)
                if hasattr(shape, 'image'):
                    try:
                        ext = (shape.image.ext or 'png').lower()
                        if ext not in ('jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp'):
                            ext = 'png'
                        b64 = base64.b64encode(shape.image.blob).decode()
                        s['imagens'].append(f'data:image/{ext};base64,{b64}')
                    except Exception:
                        pass
            except Exception:
                pass

        slides.append(s)

    return json.dumps({'total': len(slides), 'slides': slides}, ensure_ascii=False)


@bp_main.route('/apresentacao/importar', methods=['POST'])
@login_required
def importar_apresentacao():
    """Importa um PPT/PPTX como item da árvore lateral."""
    arquivo = request.files.get('arquivo_pptx')
    id_pai_raw = request.form.get('id_pai_pptx', '').strip()
    id_pai = int(id_pai_raw) if id_pai_raw.isdigit() else None

    if id_pai and not _eh_container(_topico_do_usuario(id_pai)):
        flash('A pasta de destino não foi encontrada.', 'erro')
        return redirect(url_for('main.dashboard'))

    if not arquivo or not arquivo.filename:
        flash('Escolha um arquivo PPT ou PPTX.', 'erro')
        return redirect(url_for('main.dashboard'))

    nome_seguro = secure_filename(arquivo.filename)
    ext_lower = os.path.splitext(nome_seguro)[1].lower()
    if ext_lower not in ('.ppt', '.pptx'):
        flash('Somente arquivos PPT e PPTX são suportados.', 'erro')
        return redirect(url_for('main.dashboard'))

    if ext_lower == '.ppt':
        flash('Arquivos no formato antigo .PPT não são suportados. Converta para .PPTX no PowerPoint e tente novamente.', 'erro')
        return redirect(url_for('main.dashboard'))

    conteudo_bytes = arquivo.read()
    titulo = os.path.splitext(nome_seguro)[0] or 'Apresentação'

    extracao_ok = True
    try:
        slides_json = _extrair_slides_pptx(conteudo_bytes)
    except Exception as exc:
        current_app.logger.warning('Falha ao extrair slides de "%s": %s', nome_seguro, exc)
        slides_json = '{"total":0,"slides":[]}'
        extracao_ok = False

    topico = Topico(
        id_usuario=current_user.id,
        id_pai=id_pai,
        nome=titulo,
        tipo='apresentacao',
        icone='📊',
        ordem=_proxima_ordem(id_pai),
    )
    db.session.add(topico)
    db.session.flush()
    db.session.add(Apresentacao(
        id_topico=topico.id,
        titulo=titulo,
        conteudo=conteudo_bytes,
        slides_json=slides_json,
    ))
    db.session.commit()

    if extracao_ok:
        flash('Apresentação importada com sucesso.', 'sucesso')
    else:
        flash('Arquivo salvo, mas os slides não puderam ser extraídos. Use o botão "Reprocessar" na tela da apresentação.', 'aviso')
    return redirect(url_for('main.dashboard'))


@bp_main.route('/apresentacao/<int:id>/reprocessar', methods=['POST'])
@login_required
def reprocessar_apresentacao(id):
    """Re-extrai slides do PPTX já armazenado no banco."""
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'apresentacao' or not topico.apresentacoes:
        return jsonify({'erro': 'não encontrado'}), 404
    apres = topico.apresentacoes
    if not apres.conteudo:
        return jsonify({'erro': 'arquivo não disponível'}), 404
    try:
        apres.slides_json = _extrair_slides_pptx(apres.conteudo)
        db.session.commit()
        import json
        total = json.loads(apres.slides_json).get('total', 0)
        return jsonify({'ok': True, 'total': total})
    except Exception as exc:
        current_app.logger.warning('Reprocessar apresentacao %d: %s', id, exc)
        return jsonify({'ok': False, 'erro': str(exc)}), 500


@bp_main.route('/apresentacao/<int:id>/arquivo')
@login_required
def arquivo_apresentacao(id):
    """Serve o arquivo PPTX para download."""
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'apresentacao' or not topico.apresentacoes:
        return '', 404
    apres = topico.apresentacoes
    if not apres.conteudo:
        return '', 404
    from io import BytesIO
    mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    return send_file(
        BytesIO(apres.conteudo),
        mimetype=mimetype,
        as_attachment=True,
        download_name=f'{topico.nome}.pptx',
    )


@bp_main.route('/apresentacao/<int:id>/salvar', methods=['POST'])
@login_required
def salvar_apresentacao(id):
    """Salva o slide atual visitado pelo usuário."""
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'apresentacao' or not topico.apresentacoes:
        return jsonify({'erro': 'não encontrado'}), 404
    dados = request.get_json(silent=True) or {}
    if 'slide_atual' in dados:
        topico.apresentacoes.slide_atual = max(1, int(dados['slide_atual']))
    db.session.commit()
    return jsonify({'ok': True})


@bp_main.route('/apresentacao/<int:id>/comentarios', methods=['GET'])
@login_required
def listar_comentarios_slide(id):
    """Lista comentários de um slide (query param: slide=N)."""
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'apresentacao' or not topico.apresentacoes:
        return jsonify([])
    slide_num = request.args.get('slide', type=int)
    q = ComentarioSlide.query.filter_by(id_apresentacao=topico.apresentacoes.id)
    if slide_num is not None:
        q = q.filter_by(slide_num=slide_num)
    return jsonify([{
        'id': c.id,
        'slide_num': c.slide_num,
        'texto': c.texto,
        'pos_x': c.pos_x,
        'pos_y': c.pos_y,
        'cor': c.cor,
        'criado_em': c.criado_em.strftime('%d/%m/%Y %H:%M'),
    } for c in q.order_by(ComentarioSlide.criado_em).all()])


@bp_main.route('/apresentacao/<int:id>/comentarios', methods=['POST'])
@login_required
def criar_comentario_slide(id):
    """Cria um balão de comentário em um slide."""
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'apresentacao' or not topico.apresentacoes:
        return jsonify({'erro': 'não encontrado'}), 404
    dados = request.get_json(silent=True) or {}
    texto = (dados.get('texto') or '').strip()
    if not texto:
        return jsonify({'erro': 'Texto obrigatório'}), 400
    c = ComentarioSlide(
        id_apresentacao=topico.apresentacoes.id,
        slide_num=int(dados.get('slide_num', 1)),
        texto=texto,
        pos_x=float(dados.get('pos_x', 50.0)),
        pos_y=float(dados.get('pos_y', 50.0)),
        cor=dados.get('cor', '#fbbf24'),
    )
    db.session.add(c)
    db.session.commit()
    return jsonify({
        'id': c.id,
        'slide_num': c.slide_num,
        'texto': c.texto,
        'pos_x': c.pos_x,
        'pos_y': c.pos_y,
        'cor': c.cor,
        'criado_em': c.criado_em.strftime('%d/%m/%Y %H:%M'),
    })


@bp_main.route('/apresentacao/<int:id>/comentarios/<int:cid>', methods=['DELETE'])
@login_required
def excluir_comentario_slide(id, cid):
    """Exclui um balão de comentário."""
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'apresentacao' or not topico.apresentacoes:
        return jsonify({'erro': 'não encontrado'}), 404
    c = ComentarioSlide.query.filter_by(
        id=cid, id_apresentacao=topico.apresentacoes.id
    ).first()
    if not c:
        return jsonify({'erro': 'não encontrado'}), 404
    db.session.delete(c)
    db.session.commit()
    return jsonify({'ok': True})


# ── Visualização e edição de conteúdo ─────────────────────────────────────────

_TMPL_CONTEUDO = {
    'nota':     'content/_nota.html',
    'tarefa':   'content/_nota.html',
    'caderno':  'content/_caderno.html',
    'pdf':      'content/_pdf.html',
    'imagem':   'content/_imagem.html',
    'planilha': 'content/_planilha.html',
    'jupyter':  'content/_jupyter.html',
    'diagrama': 'content/_diagrama.html',
    'python':     'content/_python.html',
    'sql':        'content/_sql.html',
    'biblioteca':   'content/_biblioteca.html',
    'apresentacao': 'content/_apresentacao.html',
}


@bp_main.route('/topico/<int:id>')
@login_required
def ver_topico(id):
    topico = _topico_do_usuario(id)
    if not topico:
        return '', 404

    tmpl = _TMPL_CONTEUDO.get(topico.tipo)
    if not tmpl:
        return '', 404

    if topico.tipo in ('nota', 'tarefa'):
        conteudo = topico.notas
    elif topico.tipo == 'caderno':
        conteudo = topico.cadernos
    elif topico.tipo == 'pdf':
        conteudo = topico.pdfs
    elif topico.tipo == 'imagem':
        conteudo = topico.imagens
    elif topico.tipo == 'planilha':
        conteudo = topico.planilhas
    elif topico.tipo == 'diagrama':
        conteudo = topico.diagramas
    elif topico.tipo == 'jupyter':
        conteudo = sorted(topico.celulas_jupyter, key=lambda c: c.ordem)
    elif topico.tipo in ('python', 'sql'):
        conteudo = topico.codigos
    elif topico.tipo == 'biblioteca':
        conteudo = (
            Topico.query
            .filter_by(id_usuario=current_user.id, id_pai=topico.id, tipo='pdf')
            .order_by(Topico.ordem, Topico.nome)
            .all()
        )
    elif topico.tipo == 'apresentacao':
        conteudo = topico.apresentacoes
    else:
        conteudo = None

    return render_template(tmpl, topico=topico, conteudo=conteudo)


@bp_main.route('/topicos/<int:id>/renomear', methods=['POST'])
@login_required
def renomear_topico(id):
    topico = _topico_do_usuario(id)
    if not topico:
        return jsonify({'erro': 'Não encontrado'}), 404
    dados = request.get_json(silent=True) or {}
    nome = dados.get('nome', '').strip()
    if not nome:
        return jsonify({'erro': 'Nome inválido'}), 400
    topico.nome = nome
    db.session.commit()
    return jsonify({'ok': True})


@bp_main.route('/topicos/<int:id>/deletar', methods=['POST'])
@login_required
def deletar_topico(id):
    topico = _topico_do_usuario(id)
    if not topico:
        return jsonify({'erro': 'Não encontrado'}), 404
    arquivos_para_remover = _coletar_arquivos_do_topico(topico)
    try:
        db.session.delete(topico)
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        return jsonify({'erro': str(e)}), 500

    _remover_arquivos_static(arquivos_para_remover)
    return jsonify({'ok': True})


@bp_main.route('/topicos/<int:id>/mover', methods=['POST'])
@login_required
def mover_topico(id):
    topico = _topico_do_usuario(id)
    if not topico:
        return jsonify({'erro': 'Não encontrado'}), 404

    dados = request.get_json(silent=True) or {}
    novo_pai_raw = dados.get('id_pai')
    novo_pai_id = int(novo_pai_raw) if novo_pai_raw is not None else None

    if novo_pai_id is not None:
        pai = _topico_do_usuario(novo_pai_id)
        if not _eh_container(pai):
            return jsonify({'erro': 'Destino deve ser uma pasta ou biblioteca'}), 400
        if novo_pai_id == topico.id:
            return jsonify({'erro': 'Não é possível mover para si mesmo'}), 400
        if _eh_descendente(topico.id, novo_pai_id):
            return jsonify({'erro': 'Não é possível mover para uma subpasta'}), 400

    topico.id_pai = novo_pai_id
    topico.ordem = _proxima_ordem(novo_pai_id)
    db.session.commit()
    return jsonify({'ok': True})


@bp_main.route('/topicos/<int:id>/reordenar', methods=['POST'])
@login_required
def reordenar_topico(id):
    """Reposiciona um item antes ou depois de um irmão na mesma pasta."""
    topico = _topico_do_usuario(id)
    if not topico:
        return jsonify({'erro': 'Não encontrado'}), 404

    dados    = request.get_json(silent=True) or {}
    alvo_id  = dados.get('alvo_id')
    posicao  = dados.get('posicao')  # 'antes' | 'depois'

    if not alvo_id or posicao not in ('antes', 'depois'):
        return jsonify({'erro': 'Parâmetros inválidos'}), 400

    alvo = _topico_do_usuario(int(alvo_id))
    if not alvo:
        return jsonify({'erro': 'Item de referência não encontrado'}), 404

    # Garante que o item vai para o mesmo pai do alvo
    topico.id_pai = alvo.id_pai

    # Busca todos os irmãos (exceto o próprio item) ordenados
    irmaos = (
        Topico.query
        .filter_by(id_usuario=current_user.id, id_pai=alvo.id_pai)
        .filter(Topico.id != topico.id)
        .order_by(Topico.ordem, Topico.id)
        .all()
    )

    # Insere o item na posição correta
    idx_alvo = next((i for i, t in enumerate(irmaos) if t.id == alvo.id), None)
    if idx_alvo is None:
        return jsonify({'erro': 'Referência não encontrada nos irmãos'}), 400

    if posicao == 'antes':
        irmaos.insert(idx_alvo, topico)
    else:
        irmaos.insert(idx_alvo + 1, topico)

    # Reatribui ordens consecutivas
    for nova_ordem, t in enumerate(irmaos, start=1):
        t.ordem = nova_ordem

    db.session.commit()
    return jsonify({'ok': True})


def _eh_descendente(topico_id, candidato_pai_id):
    """Retorna True se candidato_pai_id está na subárvore de topico_id (detecta ciclo)."""
    current_id = candidato_pai_id
    visitados = set()
    while current_id:
        if current_id == topico_id:
            return True
        if current_id in visitados:
            break
        visitados.add(current_id)
        current_id = db.session.query(Topico.id_pai).filter_by(
            id=current_id, id_usuario=current_user.id
        ).scalar()
    return False


@bp_main.route('/topico/<int:id>/salvar', methods=['POST'])
@login_required
def salvar_topico(id):
    topico = _topico_do_usuario(id)
    if not topico:
        return jsonify({'erro': 'não encontrado'}), 404

    dados = request.get_json(silent=True) or {}
    try:
        if topico.tipo in ('nota', 'tarefa') and topico.notas:
            topico.notas.conteudo = dados.get('conteudo', topico.notas.conteudo)
        elif topico.tipo == 'caderno' and topico.cadernos:
            topico.cadernos.conteudo = dados.get('conteudo', topico.cadernos.conteudo)
        elif topico.tipo == 'diagrama' and topico.diagramas:
            topico.diagramas.dados_json = dados.get('conteudo', topico.diagramas.dados_json)
        elif topico.tipo == 'planilha' and topico.planilhas:
            topico.planilhas.dados_json = dados.get('dados_json', topico.planilhas.dados_json)
        elif topico.tipo in ('python', 'sql') and topico.codigos:
            topico.codigos.conteudo = dados.get('conteudo', topico.codigos.conteudo)
        elif topico.tipo == 'pdf' and topico.pdfs:
            if 'percentual' in dados:
                topico.pdfs.percentual_lido = max(0, min(100, int(dados['percentual'])))
            if 'anotacoes' in dados:
                topico.pdfs.anotacoes = dados['anotacoes']
        elif topico.tipo == 'jupyter':
            for cel_data in dados.get('celulas', []):
                cel = CelulaJupyter.query.filter_by(
                    id=cel_data.get('id'), id_topico=topico.id
                ).first()
                if cel:
                    cel.conteudo = cel_data.get('conteudo', cel.conteudo)
        db.session.commit()
        return jsonify({'ok': True})
    except Exception as e:
        db.session.rollback()
        return jsonify({'erro': str(e)}), 500


_TIPOS_ICONE_VALIDOS = frozenset((
    'nota', 'caderno', 'planilha', 'diagrama', 'jupyter',
    'tarefa', 'imagem', 'python', 'sql', 'biblioteca', 'pdf', 'pasta', 'apresentacao',
))


@bp_main.route('/icones/upload/<tipo>', methods=['POST'])
@login_required
def upload_icone(tipo):
    from PIL import Image
    from io import BytesIO

    if tipo not in _TIPOS_ICONE_VALIDOS:
        return jsonify({'erro': 'Tipo inválido'}), 400

    arquivo = request.files.get('icone')
    if not arquivo or not arquivo.filename:
        return jsonify({'erro': 'Arquivo inválido'}), 400

    ext = arquivo.filename.rsplit('.', 1)[-1].lower()
    if ext not in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
        return jsonify({'erro': 'Formato inválido'}), 400

    pasta = os.path.join(current_app.config['UPLOAD_FOLDER'], 'icones', str(current_user.id))
    os.makedirs(pasta, exist_ok=True)

    caminho_abs = os.path.join(pasta, f'{tipo}.png')
    img = Image.open(BytesIO(arquivo.read())).convert('RGBA')
    img.thumbnail((48, 48), Image.LANCZOS)
    result = Image.new('RGBA', (48, 48), (0, 0, 0, 0))
    offset = ((48 - img.width) // 2, (48 - img.height) // 2)
    result.paste(img, offset)
    result.save(caminho_abs, 'PNG')

    caminho_rel = os.path.relpath(caminho_abs, current_app.static_folder).replace('\\', '/')
    ic = IconeCustomizado.query.filter_by(id_usuario=current_user.id, tipo=tipo).first()
    if ic:
        ic.caminho = caminho_rel
    else:
        db.session.add(IconeCustomizado(id_usuario=current_user.id, tipo=tipo, caminho=caminho_rel))
    db.session.commit()

    return jsonify({'ok': True, 'url': f'/static/{caminho_rel}'})


@bp_main.route('/topico/<int:id>/upload-capa', methods=['POST'])
@login_required
def upload_capa_pdf(id):
    from PIL import Image
    from io import BytesIO as BIO

    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'pdf' or not topico.pdfs:
        return jsonify({'erro': 'Não encontrado'}), 404

    arquivo = request.files.get('capa')
    if not arquivo or not arquivo.filename:
        return jsonify({'erro': 'Arquivo inválido'}), 400

    ext = arquivo.filename.rsplit('.', 1)[-1].lower()
    if ext not in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
        return jsonify({'erro': 'Formato inválido'}), 400

    img = Image.open(BIO(arquivo.read())).convert('RGB')

    tw, th = 200, 280
    img_aspect = img.width / img.height
    if img_aspect > tw / th:
        nw = int(img.height * tw / th)
        img = img.crop(((img.width - nw) // 2, 0, (img.width - nw) // 2 + nw, img.height))
    else:
        nh = int(img.width * th / tw)
        img = img.crop((0, (img.height - nh) // 2, img.width, (img.height - nh) // 2 + nh))
    img = img.resize((tw, th), Image.LANCZOS)

    buf = BIO()
    img.save(buf, 'JPEG', quality=85)
    topico.pdfs.conteudo_capa = buf.getvalue()
    db.session.commit()

    return jsonify({'ok': True, 'url': f'/pdf/{id}/capa'})


@bp_main.route('/pdf/<int:id>/arquivo')
@login_required
def servir_pdf(id):
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'pdf' or not topico.pdfs:
        return ('Not Found', 404)
    pdf = topico.pdfs
    if pdf.conteudo_pdf:
        import io
        return send_file(
            io.BytesIO(pdf.conteudo_pdf),
            mimetype='application/pdf',
            as_attachment=False,
            download_name=f'{topico.nome}.pdf',
        )
    if pdf.caminho:
        arquivo = _caminho_static_seguro(pdf.caminho)
        if arquivo and arquivo.is_file():
            return send_file(arquivo, mimetype='application/pdf')
    return ('Not Found', 404)


@bp_main.route('/pdf/<int:id>/capa')
@login_required
def servir_capa_pdf(id):
    topico = _topico_do_usuario(id)
    if not topico or topico.tipo != 'pdf' or not topico.pdfs:
        return ('Not Found', 404)
    pdf = topico.pdfs
    if pdf.conteudo_capa:
        resp = make_response(pdf.conteudo_capa)
        resp.headers['Content-Type'] = 'image/jpeg'
        resp.headers['Cache-Control'] = 'private, max-age=3600'
        return resp
    if pdf.capa:
        arquivo = _caminho_static_seguro(pdf.capa)
        if arquivo and arquivo.is_file():
            return send_file(arquivo, mimetype='image/jpeg')
    return ('Not Found', 404)


@bp_main.route('/topicos/buscar')
@login_required
def buscar_topicos():
    q = request.args.get('q', '').strip()
    if not q:
        return jsonify([])
    topicos = (
        Topico.query
        .filter(
            Topico.id_usuario == current_user.id,
            Topico.tipo != 'pasta',
            Topico.nome.ilike(f'%{q}%')
        )
        .order_by(Topico.nome)
        .limit(8)
        .all()
    )
    return jsonify([{'id': t.id, 'nome': t.nome, 'tipo': t.tipo} for t in topicos])


@bp_main.route('/topicos/pesquisar')
@login_required
def pesquisar_topicos():
    import re
    q = request.args.get('q', '').strip()
    if len(q) < 2:
        return jsonify([])

    vistos = {}

    por_nome = (
        Topico.query
        .filter(
            Topico.id_usuario == current_user.id,
            Topico.nome.ilike(f'%{q}%')
        )
        .order_by(Topico.nome)
        .limit(20)
        .all()
    )
    for t in por_nome:
        vistos[t.id] = {'id': t.id, 'nome': t.nome, 'tipo': t.tipo, 'match': 'nome'}

    for model_cls, join_col in [(Nota, Nota.id_topico), (Caderno, Caderno.id_topico)]:
        por_conteudo = (
            Topico.query
            .join(model_cls, join_col == Topico.id)
            .filter(
                Topico.id_usuario == current_user.id,
                model_cls.conteudo.ilike(f'%{q}%')
            )
            .order_by(Topico.nome)
            .limit(20)
            .all()
        )
        for t in por_conteudo:
            if t.id not in vistos:
                vistos[t.id] = {'id': t.id, 'nome': t.nome, 'tipo': t.tipo, 'match': 'conteúdo'}

    return jsonify(list(vistos.values())[:25])


@bp_main.route('/resumo')
@login_required
def resumo():
    """Painel de resumo — dados para a aba de dashboard."""
    hoje = date.today()

    # ── Filtro de período ─────────────────────────────────────────────────
    de_str  = request.args.get('de',  '').strip()
    ate_str = request.args.get('ate', '').strip()
    try:
        de = date.fromisoformat(de_str) if de_str else date(hoje.year, hoje.month, 1)
    except ValueError:
        de = date(hoje.year, hoje.month, 1)
    try:
        ate = date.fromisoformat(ate_str) if ate_str else hoje
    except ValueError:
        ate = hoje
    if de > ate:
        de, ate = ate, de

    de_dt  = datetime(de.year,  de.month,  de.day,  0,  0,  0)
    ate_dt = datetime(ate.year, ate.month, ate.day, 23, 59, 59)

    # ── Kanban ────────────────────────────────────────────────────────────
    colunas = (
        ColunaKanban.query
        .filter_by(id_usuario=current_user.id)
        .order_by(ColunaKanban.ordem)
        .all()
    )
    kanban_cols_raw = [
        {'nome': c.nome, 'cor': c.cor, 'total': len(c.cartoes)}
        for c in colunas
    ]
    max_cartoes = max((c['total'] for c in kanban_cols_raw), default=1) or 1
    n = len(kanban_cols_raw)
    kanban_cols = [
        {**c, 'pct': round(c['total'] / max_cartoes * 100), 'is_last': i == n - 1}
        for i, c in enumerate(kanban_cols_raw)
    ]
    total_cartoes      = sum(c['total'] for c in kanban_cols)
    cartoes_concluidos = kanban_cols[-1]['total'] if kanban_cols else 0

    def _conta_prio(p):
        return (
            CartaoKanban.query
            .join(ColunaKanban)
            .filter(ColunaKanban.id_usuario == current_user.id,
                    CartaoKanban.prioridade == p)
            .count()
        )

    prio = {
        'alta':  _conta_prio('alta'),
        'media': _conta_prio('media'),
        'baixa': _conta_prio('baixa'),
    }
    prio_total = sum(prio.values()) or 1

    # ── Tarefas no período ────────────────────────────────────────────────
    tarefas_periodo = (
        TarefaFacil.query
        .filter(
            TarefaFacil.id_usuario == current_user.id,
            TarefaFacil.data_tarefa >= de,
            TarefaFacil.data_tarefa <= ate,
        )
        .order_by(TarefaFacil.concluida, TarefaFacil.criado_em)
        .all()
    )
    concluidas_periodo = sum(1 for t in tarefas_periodo if t.concluida)
    atrasadas = (
        TarefaFacil.query
        .filter(
            TarefaFacil.id_usuario == current_user.id,
            TarefaFacil.data_tarefa < hoje,
            TarefaFacil.concluida.is_(False),
        )
        .count()
    )

    # ── Atividade no período ──────────────────────────────────────────────
    hist_rows = (
        db.session.query(
            HistoricoEtapas.nome_coluna,
            HistoricoEtapas.data_entrada,
            CartaoKanban.titulo.label('cartao_titulo'),
        )
        .join(CartaoKanban, HistoricoEtapas.id_cartao == CartaoKanban.id)
        .join(ColunaKanban, CartaoKanban.id_coluna == ColunaKanban.id)
        .filter(
            ColunaKanban.id_usuario == current_user.id,
            HistoricoEtapas.data_entrada >= de_dt,
            HistoricoEtapas.data_entrada <= ate_dt,
        )
        .order_by(desc(HistoricoEtapas.data_entrada))
        .limit(10)
        .all()
    )
    historico = [
        {'titulo': r.cartao_titulo, 'coluna': r.nome_coluna, 'data': r.data_entrada}
        for r in hist_rows
    ]

    # ── Documentos por tipo ───────────────────────────────────────────────
    _TIPO_INFO = {
        'nota':     ('📒', 'Notas',      '#2563eb'),
        'caderno':  ('📔', 'Cadernos',   '#7c3aed'),
        'planilha': ('📊', 'Planilhas',  '#16a34a'),
        'diagrama': ('🔷', 'Diagramas',  '#0891b2'),
        'jupyter':  ('⚗️',  'Notebooks',  '#d97706'),
        'tarefa':   ('🎯', 'Tarefas',    '#16a34a'),
        'imagem':   ('🖼',  'Imagens',    '#7c3aed'),
        'pdf':      ('📄', 'PDFs',       '#ef4444'),
        'python':   ('PY', 'Python',     '#3776ab'),
        'sql':      ('SQL','SQL',        '#e38c00'),
    }
    contagem: dict[str, int] = {}
    for t in Topico.query.filter_by(id_usuario=current_user.id).all():
        if t.tipo != 'pasta':
            contagem[t.tipo] = contagem.get(t.tipo, 0) + 1

    docs_info = [
        {
            'tipo':   tipo,
            'count':  count,
            'icone':  _TIPO_INFO.get(tipo, ('📄', tipo, '#64748b'))[0],
            'label':  _TIPO_INFO.get(tipo, ('📄', tipo, '#64748b'))[1],
            'cor':    _TIPO_INFO.get(tipo, ('📄', tipo, '#64748b'))[2],
        }
        for tipo, count in sorted(contagem.items(), key=lambda x: x[1], reverse=True)
    ]
    total_docs = sum(contagem.values())

    return render_template(
        'content/_resumo.html',
        kanban_cols=kanban_cols,
        total_cartoes=total_cartoes,
        cartoes_concluidos=cartoes_concluidos,
        prio=prio,
        prio_total=prio_total,
        tarefas_periodo=tarefas_periodo,
        concluidas_periodo=concluidas_periodo,
        atrasadas=atrasadas,
        historico=historico,
        docs_info=docs_info,
        total_docs=total_docs,
        hoje=hoje,
        de=de,
        ate=ate,
    )
